from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # blank slide
slide = prs.slides.add_slide(slide_layout)

# ----------------------------
# TITLE
# ----------------------------
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_tf = title_box.text_frame
title_run = title_tf.paragraphs[0].add_run()
title_run.text = "UltraTech Cement | Equity Research Snapshot"
title_run.font.size = Pt(28)
title_run.font.bold = True

# ----------------------------
# COLUMN WIDTHS
# ----------------------------
col_width = Inches(3)
left_positions = [Inches(0.3), Inches(3.4), Inches(6.5)]

section_titles = ["Market Overview", "Financial Analysis", "Strategic Insights"]

content_text = [
    "• Industry growing at 8-10%\n• Demand driven by infra push\n• Pricing power improving\n• Competitive intensity moderate\n• Regional demand variation high\n• Capacity expansion ongoing",
    
    "• Revenue: ₹70,000 Cr\n• EBITDA: ₹18,000 Cr\n• Margin: 25%\n• Strong operating leverage\n• Cost pressures easing\n• ROCE improving trend\n• Debt under control",
    
    "• Focus on premium products\n• Expansion in Tier-2 markets\n• Logistics optimization key\n• Digital transformation ongoing\n• ESG focus increasing\n• Long-term demand visibility strong"
]

# ----------------------------
# CREATE 3 SECTIONS
# ----------------------------
for i in range(3):
    left = left_positions[i]

    # Section Header Box
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, Inches(1), col_width, Inches(0.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(200, 200, 200)

    header_tf = header.text_frame
    header_tf.text = section_titles[i]
    header_tf.paragraphs[0].font.bold = True

    # Content Box
    content = slide.shapes.add_textbox(left, Inches(1.6), col_width, Inches(2.5))
    tf = content.text_frame
    tf.text = content_text[i]

    for p in tf.paragraphs:
        p.font.size = Pt(12)

    # Image Placeholder (simulate consulting clutter)
    img_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, Inches(4.2), col_width, Inches(1.5)
    )
    img_box.fill.solid()
    img_box.fill.fore_color.rgb = RGBColor(230, 230, 230)

    img_tf = img_box.text_frame
    img_tf.text = "Chart / Graph Placeholder"
    img_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ----------------------------
# ADD EXTRA CLUTTER (BOTTOM STRIP)
# ----------------------------
bottom_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(9), Inches(1))
tf_bottom = bottom_box.text_frame
tf_bottom.text = (
    "Key Takeaways:\n"
    "• Strong demand outlook supported by infra spending\n"
    "• Margin expansion driven by cost optimization\n"
    "• Valuation attractive with ~15% upside potential\n"
)

# Save file
prs.save("Cluttered_Consulting_Slide.pptx")

print("Cluttered consulting slide created!")